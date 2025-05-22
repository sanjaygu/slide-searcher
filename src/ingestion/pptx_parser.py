from pptx import Presentation
from typing import List, Dict, Any
import os
from PIL import Image
import io

class PPTXParser:
    def __init__(self):
        self.supported_extensions = ['.pptx']

    def parse(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Parse a PPTX file and extract text and images from each slide.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        presentation = Presentation(file_path)
        slides_data = []

        for idx, slide in enumerate(presentation.slides, 1):
            slide_data = {
                'slide_number': idx,
                'text_content': self._extract_text(slide),
                'images': self._extract_images(slide),
                'notes': self._extract_notes(slide),
                'layout': self._get_slide_layout(slide)
            }
            slides_data.append(slide_data)

        return slides_data

    def _extract_text(self, slide) -> str:
        """Extract text from slide shapes."""
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
        return "\n".join(text_content)

    def _extract_images(self, slide) -> List[Dict[str, Any]]:
        """Extract images from slide."""
        images = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                image_bytes = image.blob
                image_data = {
                    'position': (shape.left, shape.top),
                    'size': (shape.width, shape.height),
                    'image_bytes': image_bytes
                }
                images.append(image_data)
        return images

    def _extract_notes(self, slide) -> str:
        """Extract notes from slide."""
        if slide.has_notes_slide:
            return slide.notes_slide.notes_text_frame.text
        return ""

    def _get_slide_layout(self, slide) -> str:
        """Get the layout type of the slide."""
        return slide.slide_layout.name